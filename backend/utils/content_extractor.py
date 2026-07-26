"""
Content extraction utilities for PowerPoint presentations.

Extracts the text content of each slide of a .pptx file, in slide order,
so it can be fed to the text-to-speech step.
"""

import logging
from typing import List

from pptx import Presentation
from tqdm import tqdm

logger = logging.getLogger(__name__)


def extract_ppt_content(ppt_path: str) -> List[str]:
    """
    Extracts text content from each slide of a PowerPoint presentation.

    Args:
        ppt_path: Path to the PowerPoint (.pptx) file.

    Returns:
        A list where each element is the concatenated text of one slide's
        shapes, in slide order.

    Raises:
        FileNotFoundError: If ppt_path does not exist.
        Exception: If the file cannot be parsed as a valid .pptx.
    """
    logger.info("Loading presentation: %s", ppt_path)
    presentation = Presentation(ppt_path)
    slide_contents: List[str] = []

    total_slides = len(presentation.slides)
    logger.info("Found %d slide(s)", total_slides)

    for slide in tqdm(
        presentation.slides,
        desc="Extracting slide content",
        total=total_slides,
        unit="slide",
    ):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        slide_text.append(paragraph.text)
        slide_contents.append(" ".join(slide_text))

    return slide_contents
